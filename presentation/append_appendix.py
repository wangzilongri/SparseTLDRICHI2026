#!/usr/bin/env python3
"""
append_appendix.py — Append all appendix slides to the main talk.

Source:  output/transfer-learning-meta-analysis-ichi-2026-talk.with-05b.pptx  (13 slides)
Append:  output/appendix-lecture-method-foundations.pptx                       (30 slides)
Output:  output/transfer-learning-meta-analysis-ichi-2026-talk.full.pptx       (43 slides)

The appendix has no embedded images (pure shapes + text), so a direct
shape-tree copy is sufficient — no media relationship remapping needed.
"""

import copy, os
from pptx import Presentation
from pptx.oxml.ns import qn

BASE   = os.path.dirname(os.path.abspath(__file__))
MAIN   = os.path.join(BASE, 'output',
         'transfer-learning-meta-analysis-ichi-2026-talk.with-05b.pptx')
APPEN  = os.path.join(BASE, 'output',
         'appendix-lecture-method-foundations.pptx')
OUT    = os.path.join(BASE, 'output',
         'transfer-learning-meta-analysis-ichi-2026-talk.full.pptx')

main_prs = Presentation(MAIN)
app_prs  = Presentation(APPEN)


def append_slide(target_prs, src_slide):
    """Append src_slide as a new last slide in target_prs."""
    # Add a blank slide (layout doesn't matter — we overwrite everything)
    new_sl = target_prs.slides.add_slide(target_prs.slide_layouts[0])

    # ── Shape tree ─────────────────────────────────────────────────────────────
    tgt_sp = new_sl.shapes._spTree
    src_sp = src_slide.shapes._spTree

    # Remove auto-generated children (first 2 are nvGrpSpPr, grpSpPr — keep them)
    for child in list(tgt_sp)[2:]:
        tgt_sp.remove(child)
    for child in list(src_sp)[2:]:
        tgt_sp.append(copy.deepcopy(child))

    # ── Slide background (full-slide fills like the maroon divider slides) ─────
    # The maroon "cover" on divider slides is a full-size rectangle shape
    # already captured in the shape tree above.  Copy p:bg only if explicitly set.
    src_cSld = src_slide._element.find(qn('p:cSld'))
    tgt_cSld = new_sl._element.find(qn('p:cSld'))
    if src_cSld is not None and tgt_cSld is not None:
        src_bg = src_cSld.find(qn('p:bg'))
        if src_bg is not None:
            tgt_bg = tgt_cSld.find(qn('p:bg'))
            if tgt_bg is not None:
                tgt_cSld.remove(tgt_bg)
            tgt_cSld.insert(0, copy.deepcopy(src_bg))

    return new_sl


print(f'Main deck : {len(main_prs.slides):2d} slides  ({os.path.basename(MAIN)})')
print(f'Appendix  : {len(app_prs.slides):2d} slides  ({os.path.basename(APPEN)})')
print()

for i, sl in enumerate(app_prs.slides):
    titles = [s.text_frame.text.strip() for s in sl.shapes if s.has_text_frame]
    title  = next((t for t in titles if len(t) > 3), '(blank)')[:60]
    append_slide(main_prs, sl)
    print(f'  [{i+1:2d}] appended  {title!r}')

main_prs.save(OUT)
total = len(main_prs.slides)
print()
print(f'Saved → {OUT}')
print(f'Total  : {total} slides  (13 main + 30 appendix)')
