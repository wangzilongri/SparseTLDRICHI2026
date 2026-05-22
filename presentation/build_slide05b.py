#!/usr/bin/env python3
"""
build_slide05b.py — Insert candidate slide 05b (equation detail) after slide 05.

Layout: 4 horizontal rows, each split:
  Left (244px)  — step number + name + description (light bg, accent strip)
  Right (886px) — MAROON panel with white-text equation PNG + annotation

Usage:  python3 presentation/build_slide05b.py
Output: presentation/output/transfer-learning-meta-analysis-ichi-2026-talk.with-05b.pptx
"""

import io, os, struct
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── pixel → EMU (96 DPI, 1280×720 design space) ───────────────────────────────
def em(px): return Emu(int(round(px * 9525)))

def c(h):
    h = h.lstrip('#')
    return RGBColor(int(h[:2],16), int(h[2:4],16), int(h[4:],16))

MAROON  = c('#7A0019')
GOLD    = c('#D8AE4B')
BLUE    = c('#4A6F8A')
GREEN   = c('#4A6A4A')
TEXT_C  = c('#231F20')
MUTED_C = c('#5E5A56')
WHITE   = c('#FFFFFF')
LINE_C  = c('#E4DBD2')
BGROW   = c('#FAFAF8')
ANNOT_C = c('#FFD080')  # soft gold for annotation text on maroon

# ── Paths ──────────────────────────────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
SRC  = os.path.join(BASE, 'output',
       'transfer-learning-meta-analysis-ichi-2026-talk.results-arc.pptx')
OUT  = os.path.join(BASE, 'output',
       'transfer-learning-meta-analysis-ichi-2026-talk.with-05b.pptx')
EQ   = os.path.join(BASE, 'assets', 'equations')

# ── Load ───────────────────────────────────────────────────────────────────────
prs = Presentation(SRC)

# Extract logo bytes from slide 05 (index 4) before any reordering
logo_blob = None
for shp in prs.slides[4].shapes:
    if shp.shape_type == 13:          # MSO_SHAPE_TYPE.PICTURE
        logo_blob = shp.image.blob
        break

# ── Add blank slide at end, then move to position 5 ───────────────────────────
new_sl = prs.slides.add_slide(prs.slide_layouts[0])

from pptx.oxml.ns import qn as _qn
slist = prs.part._element.find(_qn('p:sldIdLst'))
ids   = list(slist)
last  = ids[-1]
slist.remove(last)
slist.insert(5, last)          # right after slide 05 (index 4)

sh = new_sl.shapes

# ── Shape helpers ──────────────────────────────────────────────────────────────
def add_rect(x, y, w, h, fill, line=None, lw=0.75, rounded=False):
    s = sh.add_shape(5 if rounded else 1, em(x), em(y), em(w), em(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        try:
            s.line.fill.background()
        except Exception:
            s.line.color.rgb = fill   # invisible fallback
    return s

def add_text(x, y, w, h, text, sz=13, col=TEXT_C, bold=False,
             align=PP_ALIGN.LEFT, face='Aptos', wrap=True):
    tb = sh.add_textbox(em(x), em(y), em(w), em(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.color.rgb = col
    r.font.bold = bold
    r.font.name = face
    return tb

def add_pic(img, x, y, w, h):
    return sh.add_picture(img, em(x), em(y), em(w), em(h))

def png_dims(path):
    with open(path, 'rb') as f:
        f.read(8)            # PNG signature
        f.read(4)            # chunk length
        f.read(4)            # IHDR marker
        w = struct.unpack('>I', f.read(4))[0]
        h = struct.unpack('>I', f.read(4))[0]
    return w, h

# ── Chrome ─────────────────────────────────────────────────────────────────────
add_rect(0, 0, 1280, 16, MAROON)
add_rect(0, 16, 1280, 3, GOLD)

if logo_blob:
    add_pic(io.BytesIO(logo_blob), 1050.68, 34, 168.64, 72)

add_text(72, 34, 940, 56,
         'Estimator pipeline — equations from §3',
         sz=29, col=MAROON, bold=True, face='Aptos Display')

add_text(1004, 48, 198, 18, '(candidate 5b)',
         sz=10, col=MUTED_C, align=PP_ALIGN.RIGHT)

add_rect(72, 670, 1136, 2, LINE_C)
add_text(72, 679, 400, 20, 'IEEE ICHI 2026', sz=12, col=MUTED_C)
add_text(840, 679, 368, 20, 'Minneapolis, MN • June 1–3, 2026',
         sz=12, col=MUTED_C, align=PP_ALIGN.RIGHT)

add_rect(1058, 626, 150, 26, WHITE, LINE_C, lw=0.75, rounded=True)
add_text(1070, 631, 126, 16, '4:15–6:00  [5b]',
         sz=11, col=MUTED_C, align=PP_ALIGN.CENTER)

# ── Row layout constants ───────────────────────────────────────────────────────
ROW_H   = 130
ROW_GAP = 6
ROW_YS  = [108, 108 + ROW_H + ROW_GAP,
               108 + 2*(ROW_H + ROW_GAP),
               108 + 3*(ROW_H + ROW_GAP)]  # [108, 244, 380, 516]

EQ_X = 322   # left edge of equation panel
EQ_W = 886   # width of equation panel

ROWS = [
    dict(
        accent=GOLD,
        num='01', title='Fit source models',
        desc='K source RCTs via glmtrans.\ne(x) known by design (A2).',
        eq='source-proxy.png',
        annot='μ̂₀⁺ (x),  μ̂₁⁺ (x)   k=1,…,K     \xa7 3.1',
    ),
    dict(
        accent=BLUE,
        num='02', title='Anchor baseline',
        desc='LASSO on target placebo\nresiduals; corrects bias.',
        eq='sparse-correction.png',
        eq2='target-anchor.png',
        annot='δ̂ → μ̂₀ᵃⁿᶜʰᵒʳ     \xa7 3.2   Eq. (4)',
    ),
    dict(
        accent=GREEN,
        num='03', title='Build pseudo-outcomes',
        desc='DR pseudo-outcome ψᵢ;\ncross-fitting (Proposed-CF).',
        eq='dr-learner.png',
        annot='ψᵢ   (Kennedy 2023)     Eq. (5)',
    ),
    dict(
        accent=BLUE,
        num='04', title='Estimate CATE',
        desc='Regress ψᵢ on covariates.\nThm 1: n₀⁻½ rate (connected).',
        eq=None,
        text_eq='τ̂(x)  =  argmin Σᵢ (ψᵢ − f(Xᵢ))²',
        text_eq2='f  ∈  {random forest, lasso, …}',
        annot='τ̂(x)     Thm. 1, 2     \xa7 3.3',
    ),
]

# ── Draw rows ──────────────────────────────────────────────────────────────────
for i, row in enumerate(ROWS):
    y   = ROW_YS[i]
    acc = row['accent']

    # Left label panel
    add_rect(72, y, 244, ROW_H, BGROW, LINE_C, lw=0.5)
    add_rect(72, y, 5, ROW_H, acc)

    # Step badge
    add_rect(83, y+8, 38, 38, BGROW, acc, lw=1.5, rounded=True)
    add_text(83, y+13, 38, 28, row['num'],
             sz=18, col=acc, bold=True, align=PP_ALIGN.CENTER)

    # Step title + description
    add_text(130, y+10, 180, 26, row['title'], sz=14, col=TEXT_C, bold=True)
    add_text(130, y+40, 180, 82, row['desc'], sz=11, col=MUTED_C, wrap=True)

    # Equation panel (MAROON background)
    add_rect(EQ_X, y, EQ_W, ROW_H, MAROON)

    if row.get('eq'):
        eq_path = os.path.join(EQ, row['eq'])
        eq_w, eq_h = png_dims(eq_path)

        if row.get('eq2'):
            # Stack two equations (row 02)
            eq2_path = os.path.join(EQ, row['eq2'])
            eq2_w, eq2_h = png_dims(eq2_path)

            d1_w = int(EQ_W * 0.72)
            d1_h = int(eq_h  * (d1_w / eq_w))
            d2_w = int(EQ_W * 0.68)
            d2_h = int(eq2_h * (d2_w / eq2_w))

            # Shrink until both fit with 8px gap and 20px margins
            while d1_h + 8 + d2_h > ROW_H - 24:
                d1_w = int(d1_w * 0.95);  d1_h = int(eq_h  * (d1_w / eq_w))
                d2_w = int(d2_w * 0.95);  d2_h = int(eq2_h * (d2_w / eq2_w))

            total = d1_h + 8 + d2_h
            y1 = y + (ROW_H - total) // 2
            y2 = y1 + d1_h + 8
            x1 = EQ_X + (EQ_W - d1_w) // 2
            x2 = EQ_X + (EQ_W - d2_w) // 2

            add_pic(eq_path,  x1, y1, d1_w, d1_h)
            add_pic(eq2_path, x2, y2, d2_w, d2_h)

        else:
            # Single equation, centered
            d_w = min(int(EQ_W * 0.75), 660)
            d_h = int(eq_h * (d_w / eq_w))
            if d_h > ROW_H - 24:
                d_h = ROW_H - 24
                d_w = int(eq_w * (d_h / eq_h))
            cx = EQ_X + (EQ_W - d_w) // 2
            cy = y  + (ROW_H - d_h) // 2
            add_pic(eq_path, cx, cy, d_w, d_h)

    else:
        # Text equation for step 04
        add_text(EQ_X+20, y+22, EQ_W-40, 38,
                 row['text_eq'],
                 sz=17, col=WHITE, face='Courier New', align=PP_ALIGN.CENTER)
        add_text(EQ_X+20, y+62, EQ_W-40, 26,
                 row['text_eq2'],
                 sz=12, col=c('#FFD0D0'), face='Courier New', align=PP_ALIGN.CENTER)

    # Annotation (bottom-right of equation panel)
    if row.get('annot'):
        add_text(EQ_X+4, y+ROW_H-22, EQ_W-10, 20,
                 row['annot'], sz=10, col=ANNOT_C, align=PP_ALIGN.RIGHT)

# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f'OK  →  {OUT}')
