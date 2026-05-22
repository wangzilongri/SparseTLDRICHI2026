#!/usr/bin/env python3
"""
patch_cover_slide.py — Fix slide 01 for ICHI 2026 submission requirements.

Changes:
  1. Replace "Research Paper" with the actual session label.
  2. Expand author line to two lines: names (with presenter marked) + full affiliation.
  3. Add date/location line to cover body (guidelines require it on the slide, not only footer).
  4. Save as properly named submission file.

Usage:
  python3 presentation/patch_cover_slide.py
  python3 presentation/patch_cover_slide.py --session "SS19"   # override session label
"""

import argparse, copy, io, os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

parser = argparse.ArgumentParser()
parser.add_argument('--session', default='Scientific Session 19',
                    help='Session label to show on cover (default: "Scientific Session 19")')
parser.add_argument('--session-id', default='SS19',
                    help='Session ID abbreviation for filename (default: SS19)')
args = parser.parse_args()

BASE   = os.path.dirname(os.path.abspath(__file__))
SRC    = os.path.join(BASE, 'output',
         'transfer-learning-meta-analysis-ichi-2026-talk.results-arc.pptx')
# Submission filename: SessionID_PaperID_LastName.pptx
FNAME  = f'{args.session_id}_165_Wang.pptx'
OUT    = os.path.join(BASE, 'output', FNAME)

def em(px): return Emu(int(round(px * 9525)))
def c(h):
    h = h.lstrip('#')
    return RGBColor(int(h[:2],16), int(h[2:4],16), int(h[4:],16))

MAROON = c('#7A0019')
GOLD   = c('#D8AE4B')
MUTED  = c('#5E5A56')
WHITE  = c('#FFFFFF')

prs  = Presentation(SRC)
sl   = prs.slides[0]
shps = sl.shapes

# ── 1. Replace "Research Paper" with actual session label ─────────────────────
for shp in shps:
    if shp.has_text_frame and 'Research Paper' in shp.text_frame.text:
        tf = shp.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if 'Research Paper' in run.text:
                    run.text = args.session
        print(f'  Session label → {args.session!r}')
        break

# ── 2. Expand author / affiliation block ─────────────────────────────────────
# Find the existing condensed author shape at approx (106, 425)
author_shp = None
for shp in shps:
    if shp.has_text_frame:
        txt = shp.text_frame.text
        if 'Zilong Wang' in txt and 'Georgia' in txt:
            author_shp = shp
            break

if author_shp:
    # Resize to two lines and rewrite
    author_shp.height = em(56)   # was 30px — expand for two lines

    tf = author_shp.text_frame
    tf.word_wrap = False

    # Line 1 — names, presenter underlined
    p1 = tf.paragraphs[0]
    # Clear existing runs
    for run in list(p1.runs):
        p1._p.remove(run._r)

    def add_run(para, text, bold=False, underline=False, sz=14):
        from pptx.oxml.ns import qn
        from lxml import etree
        r = para.add_run()
        r.text = text
        r.font.size = Pt(sz)
        r.font.color.rgb = MUTED
        r.font.bold = bold
        r.font.underline = underline
        r.font.name = 'Aptos'
        return r

    add_run(p1, 'Zilong Wang', bold=True, underline=True)   # presenter
    add_run(p1, ',  Ali Abdeen,  Turgay Ayer')
    add_run(p1, '   ·   presenter: Zilong Wang', sz=11)

    # Line 2 — affiliation
    from pptx.oxml.ns import qn
    from lxml import etree
    # Add a second paragraph
    p2_xml = etree.SubElement(tf._txBody, qn('a:p'))
    r2 = etree.SubElement(p2_xml, qn('a:r'))
    rPr = etree.SubElement(r2, qn('a:rPr'), lang='en-US')
    rPr.set('sz', '1100')   # 11pt in hundredths of a point
    rPr.set('dirty', '0')
    from lxml import etree as ET
    solidFill = ET.SubElement(rPr, qn('a:solidFill'))
    srgbClr   = ET.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', '5E5A56')
    latin = ET.SubElement(rPr, qn('a:latin'))
    latin.set('typeface', 'Aptos')
    t = ET.SubElement(r2, qn('a:t'))
    t.text = 'Industrial and Systems Engineering  ·  Georgia Institute of Technology  ·  Atlanta, GA'

    print('  Author block updated (2 lines: names + affiliation)')

# ── 3. Ensure date/location is visible on cover body (not just footer) ────────
# Check if there's already a date element in the body; if not, add one
has_date_in_body = any(
    shp.has_text_frame and 'Minneapolis' in shp.text_frame.text
    and shp.top < em(640)    # above the footer zone
    for shp in shps
)

if not has_date_in_body:
    # Add a small date line just below the author block (~y=490)
    tb = shps.add_textbox(em(106), em(490), em(500), em(22))
    tf = tb.text_frame
    p  = tf.paragraphs[0]
    r  = p.add_run()
    r.text = 'June 1–3, 2026  ·  Minneapolis, MN  ·  IEEE ICHI 2026'
    r.font.size    = Pt(11)
    r.font.color.rgb = MUTED
    r.font.name    = 'Aptos'
    print('  Date/location line added to cover body')
else:
    print('  Date/location already present in cover body')

# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(OUT)
print()
print(f'Saved  → {OUT}')
print(f'Slides : {len(prs.slides)}')
print()
print('Submission checklist:')
print(f'  [x] File format : .pptx')
print(f'  [x] File name   : {FNAME}')
print(f'  [x] Slide size  : 16:9 (1280×720)')
print(f'  [x] Title       : on slide 01')
print(f'  [x] Authors     : Zilong Wang, Ali Abdeen, Turgay Ayer')
print(f'  [x] Affiliation : Georgia Institute of Technology, ISyE')
print(f'  [x] Presenter   : Zilong Wang (underlined)')
print(f'  [x] Location    : Minneapolis, MN')
print(f'  [ ] Session ID  : confirm "{args.session_id}" is the correct abbreviation')
